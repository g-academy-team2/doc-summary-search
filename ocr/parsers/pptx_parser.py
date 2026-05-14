import logging
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from .vision_helper import get_text_from_image_bytes

logger = logging.getLogger(__name__)

def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        full_text = []
        
        for i, slide in enumerate(prs.slides):
            logger.info(f"🎞️ {i+1}번 슬라이드 분석 중...")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    full_text.append(shape.text)
                
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_text:
                            full_text.append(" | ".join(row_text))

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_bytes = shape.image.blob
                        vision_text = get_text_from_image_bytes(image_bytes)
                        if vision_text:
                            full_text.append(f"\n[슬라이드 이미지 OCR]:\n{vision_text}")
                    except Exception as img_err:
                        logger.warning(f"⚠️ 이미지 분석 건너뜀: {img_err}")

            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    full_text.append(f"[슬라이드 노트]: {notes}")
                    
        return "\n\n".join(full_text)

    except Exception as e:
        logger.error(f"❌ PPTX 파싱 에러 ({file_path}): {e}")
        return f"PPTX 읽기 실패: {str(e)}"