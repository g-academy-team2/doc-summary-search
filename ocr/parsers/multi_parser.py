import os
import re # ✨ 데이터 세탁을 위한 정규표현식 도구 추가!
from docx import Document
from pptx import Presentation
import olefile
from pdf_parser import extract_text_from_pdf

# === [데이터 엔지니어의 마법: 텍스트 세탁기] ===
def clean_text(text):
    if not text:
        return "❌ 텍스트를 추출하지 못했습니다."
    
    # 1. 불필요한 특수문자, 깨진 기호 청소 (알파벳, 한글, 숫자, 기본 문장부호만 남김)
    text = re.sub(r'[^\w\s\.,!?가-힣a-zA-Z0-9\-\(\)\[\]]', ' ', text)
    
    # 2. 쓸데없이 길게 이어진 띄어쓰기(스페이스바)를 딱 한 칸으로 압축
    text = re.sub(r' {2,}', ' ', text)
    
    # 3. 의미 없이 뚝뚝 끊긴 줄바꿈을 자연스럽게 이어줌 (단, 문단 구분을 위해 연속된 줄바꿈은 2개까지만 허용)
    text = re.sub(r'\n{3,}', '\n\n', text)
    
    # 4. 앞뒤에 남은 지저분한 여백 가위표!
    return text.strip()
# ================================================

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def extract_text_from_hwp(file_path):
    f = olefile.OleFileIO(file_path)
    dirs = f.listdir()
    if ["PrvText"] in dirs:
        data = f.openstream("PrvText").read()
        return data.decode("utf-16")
    else:
        return "❌ HWP 본문 스트림을 찾을 수 없습니다."

def total_parser(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    print(f"📂 [{ext}] 형식 문서를 분석하고 세탁기를 가동합니다...")

    raw_text = ""
    if ext == ".pdf":
        raw_text = extract_text_from_pdf(file_path)
    elif ext == ".docx":
        raw_text = extract_text_from_docx(file_path)
    elif ext == ".pptx":
        raw_text = extract_text_from_pptx(file_path)
    elif ext == ".hwp":
        raw_text = extract_text_from_hwp(file_path)
    else:
        return "❌ 지원하지 않는 파일 형식입니다."
    
    # 추출된 날것의 텍스트를 세탁기에 통과시켜서 반환!
    return clean_text(raw_text)

if __name__ == "__main__":
    # 아까 테스트했던 파워포인트나 한글 파일을 그대로 다시 돌려봅니다!
    test_file = "sample.hwp" 
    
    if os.path.exists(test_file):
        print(total_parser(test_file)[:1000]) # 1000자까지만 미리보기
    else:
        print(f"⚠️ 테스트할 {test_file} 파일이 폴더에 없습니다.")