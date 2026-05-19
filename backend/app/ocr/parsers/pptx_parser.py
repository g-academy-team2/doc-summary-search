import os
import logging
import uuid # 💡 임시 파일 이름을 위해 추가
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 💡 1. 새로운 비전 헬퍼 함수명과 절대 경로로 수정했습니다.
from app.ocr.parsers.vision_helper import extract_text_from_image

logger = logging.getLogger(__name__)

def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        full_text = []
        
        for i, slide in enumerate(prs.slides):
            logger.info(f"🎞️ {i+1}번 슬라이드 분석 중...")
            
            for shape in slide.shapes:
                # 텍스트 상자 추출
                if hasattr(shape, "text") and shape.text.strip():
                    full_text.append(shape.text)
                
                # 표(Table) 추출
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_text:
                            full_text.append(" | ".join(row_text))

                # 💡 이미지 처리 (임시 파일 생성 방식 적용)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_bytes = shape.image.blob
                        
                        # 2. 바이트 데이터를 임시 이미지 파일로 저장
                        temp_img_path = f"temp_pptx_img_{uuid.uuid4().hex}.png"
                        with open(temp_img_path, "wb") as f:
                            f.write(image_bytes)
                        
                        # 3. 실제 파일 경로를 넘겨서 비전 헬퍼에게 읽게 함
                        vision_text = extract_text_from_image(temp_img_path)
                        
                        # 4. 다 쓴 임시 파일은 즉시 삭제!
                        if os.path.exists(temp_img_path):
                            os.remove(temp_img_path)
                            
                        if vision_text:
                            full_text.append(f"\n[슬라이드 이미지 OCR]:\n{vision_text}")
                            
                    except Exception as img_err:
                        logger.warning(f"⚠️ 이미지 분석 건너뜀: {img_err}")

            # 슬라이드 노트 추출
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    full_text.append(f"[슬라이드 노트]: {notes}")
                    
        return "\n\n".join(full_text)

    except Exception as e:
        logger.error(f"❌ PPTX 파싱 에러 ({file_path}): {e}")
        return f"PPTX 읽기 실패: {str(e)}"